from pptx import Presentation

# Extract powerpoint notes from Powerpoint presentation
def extract_notes(pptx_path, output_path="speaker_notes.txt"):
    prs = Presentation(pptx_path)
    with open(output_path, "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides, start=1):
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                f.write(f"Slide {i}:\n{notes_text}\n\n")
            else:
                f.write(f"Slide {i}:\n(No notes)\n\n")
    print(f"Speaker notes saved to: {output_path}")

# Example usage:
extract_notes("your_presentation.pptx")
